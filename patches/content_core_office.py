import asyncio
from functools import partial

from docx import Document  # type: ignore
from openpyxl import load_workbook  # type: ignore
from pptx import Presentation  # type: ignore

from content_core.common import ProcessSourceState
from content_core.logging import logger
from docx.table import Table
from docx.text.paragraph import Paragraph

SUPPORTED_OFFICE_TYPES = [
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
]

async def extract_docx_content_detailed(file_path):
    """Extract content from DOCX file preserving order of paragraphs and tables"""

    def _extract():
        try:
            doc = Document(file_path)
            content = []

            # We iterate through the underlying XML body to maintain exact document order
            # This ensures Tables appearing between paragraphs are caught correctly.
            for element in doc.element.body:
                # HANDLE PARAGRAPHS
                if isinstance(element, Paragraph) or hasattr(element, 'text'):
                    # Convert raw element to docx Paragraph object if needed
                    para = Paragraph(element, doc) if not isinstance(element, Paragraph) else element
                    
                    if not para.text.strip():
                        continue

                    style = para.style.name if para.style else "Normal"
                    text = para.text.strip()

                    # Indentation Logic
                    p_format = para.paragraph_format
                    indent = p_format.left_indent or 0
                    indent_level = 0
                    if hasattr(indent, "pt"):
                        indent_level = int(indent.pt / 72)
                    indent_spaces = " " * (indent_level * 4)

                    # IMPROVED HEADER DETECTION (Catches "PART-I" even if not styled as Heading)
                    is_bold_header = all(run.bold for run in para.runs if run.text.strip())
                    
                    if "Heading" in style or (is_bold_header and len(text) < 100):
                        level = style[-1] if style[-1].isdigit() else "1"
                        heading_marks = "#" * int(level)
                        content.append(f"\n{heading_marks} {text}\n")

                    # List Logic
                    elif style.startswith("List") or (hasattr(para._p, "pPr") and para._p.pPr.numPr is not None):
                        content.append(f"{indent_spaces}* {text}")

                    else:
                        # Standard Text Formatting
                        formatted_text = []
                        for run in para.runs:
                            if run.bold:
                                formatted_text.append(f"**{run.text}**")
                            elif run.italic:
                                formatted_text.append(f"*{run.text}*")
                            else:
                                formatted_text.append(run.text)
                        content.append(f"{indent_spaces}{''.join(formatted_text)}")

                # HANDLE TABLES (Ensures they are extracted where they appear)
                elif isinstance(element, Table) or (hasattr(element, 'tag') and 'tbl' in element.tag):
                    # Find the corresponding Table object
                    # For safety in some versions, we match by index if direct conversion is tricky
                    table = Table(element, doc) if not isinstance(element, Table) else element

                    rows = table.rows
                    if not rows:
                        continue

                    # Build a 2D grid to handle merged cells (cells can span multiple cols/rows)
                    # We track unique cell ids to avoid duplicating merged cell content
                    seen_cell_ids = set()
                    grid = []
                    for row in rows:
                        row_data = []
                        for cell in row.cells:
                            cell_id = id(cell._tc)
                            if cell_id in seen_cell_ids:
                                # Merged cell — represent as empty to keep column alignment
                                row_data.append("")
                            else:
                                seen_cell_ids.add(cell_id)
                                # Preserve inner newlines as spaces, strip outer whitespace
                                cell_text = cell.text.strip().replace("\n", " ").replace("|", "\\|")
                                row_data.append(cell_text)
                        if any(row_data):
                            grid.append(row_data)

                    if grid:
                        # Pad all rows to the same column count
                        max_cols = max(len(r) for r in grid)
                        for r in grid:
                            while len(r) < max_cols:
                                r.append("")

                        # Emit as a proper markdown table
                        # First row is treated as the header
                        header = grid[0]
                        separator = ["---"] * max_cols
                        content.append("\n| " + " | ".join(header) + " |")
                        content.append("| " + " | ".join(separator) + " |")
                        for data_row in grid[1:]:
                            content.append("| " + " | ".join(data_row) + " |")
                        content.append("")  # blank line after table

            return "\n\n".join(content)

        except Exception as e:
            # Using print as a fallback for logger
            print(f"Failed to extract DOCX content: {e}")
            return None

    return await asyncio.get_event_loop().run_in_executor(None, _extract)


async def get_docx_info(file_path):
    """Get DOCX metadata and content"""

    async def _get_info():
        try:
            doc = Document(file_path)

            # Extract core properties if available
            core_props = {
                "author": doc.core_properties.author,
                "created": doc.core_properties.created,
                "modified": doc.core_properties.modified,
                "title": doc.core_properties.title,
                "subject": doc.core_properties.subject,
                "keywords": doc.core_properties.keywords,
                "category": doc.core_properties.category,
                "comments": doc.core_properties.comments,
            }

            # Get document content
            content = await extract_docx_content_detailed(file_path)

            # Get document statistics
            stats = {
                "paragraph_count": len(doc.paragraphs),
                "word_count": sum(
                    len(p.text.split()) for p in doc.paragraphs if p.text.strip()
                ),
                "character_count": sum(
                    len(p.text) for p in doc.paragraphs if p.text.strip()
                ),
            }

            return {"metadata": core_props, "content": content, "statistics": stats}

        except Exception as e:
            logger.error(f"Failed to get DOCX info: {e}")
            return None

    return await _get_info()


async def extract_pptx_content(file_path):
    """Extract content from PPTX file"""

    def _extract():
        try:
            prs = Presentation(file_path)
            content = []

            for slide_number, slide in enumerate(prs.slides, 1):
                content.append(f"\n# Slide {slide_number}\n")

                # Extract title
                if slide.shapes.title:
                    content.append(f"## {slide.shapes.title.text}\n")

                # Extract text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        if (
                            shape != slide.shapes.title
                        ):  # Skip title as it's already added
                            content.append(shape.text.strip())

            return "\n\n".join(content)

        except Exception as e:
            logger.error(f"Failed to extract PPTX content: {e}")
            return None

    return await asyncio.get_event_loop().run_in_executor(None, _extract)


async def extract_xlsx_content(file_path, max_rows=10000, max_cols=100):
    """Extract content from XLSX file"""

    def _extract():
        try:
            wb = load_workbook(file_path, data_only=True)
            content = []

            for sheet in wb.sheetnames:
                ws = wb[sheet]
                content.append(f"\n# Sheet: {sheet}\n")

                # Get the maximum row and column with data
                max_row = min(ws.max_row, max_rows)
                max_col = min(ws.max_column, max_cols)

                # Create markdown table header
                headers = []
                for col in range(1, max_col + 1):
                    cell_value = ws.cell(row=1, column=col).value
                    headers.append(str(cell_value) if cell_value is not None else "")

                content.append("| " + " | ".join(headers) + " |")
                content.append("| " + " | ".join(["---"] * len(headers)) + " |")

                # Add table content
                for row in range(2, max_row + 1):
                    row_data = []
                    for col in range(1, max_col + 1):
                        cell_value = ws.cell(row=row, column=col).value
                        row_data.append(
                            str(cell_value) if cell_value is not None else ""
                        )
                    content.append("| " + " | ".join(row_data) + " |")

            return "\n".join(content)

        except Exception as e:
            logger.error(f"Failed to extract XLSX content: {e}")
            return None

    return await asyncio.get_event_loop().run_in_executor(None, partial(_extract))


async def get_pptx_info(file_path):
    """Get PPTX metadata and content"""

    def _get_pptx_metadata_sync(file_path):
        """Synchronous helper to extract metadata using python-pptx."""
        try:
            prs = Presentation(file_path)
            props = {
                "slide_count": len(prs.slides),
                "title": "",  # PowerPoint doesn't have built-in metadata like Word
            }
            stats = {
                "slide_count": len(prs.slides),
                "shape_count": sum(len(slide.shapes) for slide in prs.slides),
                "text_frame_count": sum(
                    sum(1 for shape in slide.shapes if hasattr(shape, "text"))
                    for slide in prs.slides
                ),
            }
            return {"metadata": props, "statistics": stats}
        except Exception as e:
            logger.error(f"Failed to get PPTX metadata: {e}")
            return None

    try:
        # Run blocking python-pptx operations in executor
        metadata_info = await asyncio.get_event_loop().run_in_executor(
            None, _get_pptx_metadata_sync, file_path
        )

        # Await the async content extraction directly
        content = await extract_pptx_content(file_path)

        if metadata_info:
            # Combine results
            return {**metadata_info, "content": content}
        else:
            # Fallback if metadata extraction failed
            return {"metadata": {}, "statistics": {}, "content": content}

    except Exception as e:
        logger.error(f"Failed to get PPTX info: {e}")
        return None


async def get_xlsx_info(file_path):
    """Get XLSX metadata and content"""

    async def _get_info():
        try:
            wb = load_workbook(file_path, data_only=True)

            # Extract basic properties
            props = {
                "sheet_count": len(wb.sheetnames),
                "sheets": wb.sheetnames,
                "title": wb.properties.title,
                "creator": wb.properties.creator,
                "created": wb.properties.created,
                "modified": wb.properties.modified,
            }

            # Get document content
            content = await extract_xlsx_content(file_path)

            # Get workbook statistics
            stats = {
                "sheet_count": len(wb.sheetnames),
                "total_rows": sum(sheet.max_row for sheet in wb.worksheets),
                "total_columns": sum(sheet.max_column for sheet in wb.worksheets),
            }

            return {"metadata": props, "content": content, "statistics": stats}

        except Exception as e:
            logger.error(f"Failed to get XLSX info: {e}")
            return None

    return await _get_info()


async def extract_office_content(state: ProcessSourceState):
    """Universal function to extract content from Office files"""
    assert state.file_path, "No file path provided"
    assert state.identified_type in SUPPORTED_OFFICE_TYPES, "Unsupported File Type"
    file_path = state.file_path
    doc_type = state.identified_type

    if (
        doc_type
        == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    ):
        logger.debug("Extracting content from DOCX file")
        content = await extract_docx_content_detailed(file_path)
        info = await get_docx_info(file_path)
    elif (
        doc_type
        == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    ):
        logger.debug("Extracting content from PPTX file")
        content = await extract_pptx_content(file_path)
        info = await get_pptx_info(file_path)
    elif (
        doc_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    ):
        logger.debug("Extracting content from XLSX file")
        content = await extract_xlsx_content(file_path)
        info = await get_xlsx_info(file_path)
    else:
        raise Exception(f"Unsupported file format: {doc_type}")

    del info["content"]
    return {"content": content, "metadata": info}
    