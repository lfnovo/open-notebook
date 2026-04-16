"""
Microsoft Office file processors
Handles PPT, PPTX, XLS, XLSX formats
"""

import subprocess
from pathlib import Path
from typing import Optional

from loguru import logger


def extract_text_from_ppt(ppt_path: str) -> str:
    """
    Extract text from legacy .ppt file using catppt or textract.
    
    Args:
        ppt_path: Path to the .ppt file
        
    Returns:
        Extracted text content
    """
    try:
        logger.info(f"Extracting text from .ppt file: {Path(ppt_path).name}")
        
        # First check if it's actually a PPTX file (ZIP format) misnamed as .ppt
        import zipfile
        abs_path = ppt_path if ppt_path.startswith('/') else f"/app/{ppt_path}"
        
        if zipfile.is_zipfile(abs_path):
            logger.info("File is actually PPTX format (ZIP), using python-pptx")
            try:
                from pptx import Presentation
                prs = Presentation(abs_path)
                text_runs = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_runs.append(shape.text)
                text = "\n".join(text_runs)
                
                if text.strip():
                    logger.info(f"Extracted {len(text)} characters from PPTX file using python-pptx")
                    return text
                else:
                    logger.warning("PPTX file contains no text")
                    
            except ImportError:
                logger.warning("python-pptx not installed")
            except Exception as e:
                logger.warning(f"python-pptx failed: {e}")
        
        # Try using catppt (part of catdoc package) for true legacy PPT
        try:
            result = subprocess.run(
                ['catppt', ppt_path],
                capture_output=True,
                text=True,
                check=True,
                timeout=120
            )
            text = result.stdout
            
            if text.strip():
                logger.info(f"Extracted {len(text)} characters from .ppt file using catppt")
                return text
            else:
                logger.warning("catppt returned empty output, trying textract")
                
        except FileNotFoundError:
            logger.warning("catppt not found, trying textract")
        except subprocess.CalledProcessError as e:
            logger.warning(f"catppt failed with exit code {e.returncode}: {e.stderr}, trying textract")
        
        # Fallback to textract
        try:
            import textract
            text = textract.process(ppt_path).decode('utf-8')
            
            if text.strip():
                logger.info(f"Extracted {len(text)} characters from .ppt file using textract")
                return text
                
        except ImportError:
            logger.error("textract not installed")
        except Exception as e:
            logger.error(f"textract failed: {e}")
        
        raise ValueError(
            "Unable to extract text from .ppt file. "
            "The file may be empty, corrupted, or in an unsupported format."
        )
        
    except Exception as e:
        logger.error(f"Failed to extract text from .ppt file {ppt_path}: {e}")
        raise


def extract_text_from_xls(xls_path: str) -> str:
    """
    Extract text from legacy .xls file using xls2csv or textract.
    
    Args:
        xls_path: Path to the .xls file
        
    Returns:
        Extracted text content
    """
    try:
        logger.info(f"Extracting text from .xls file: {Path(xls_path).name}")
        
        # Try using xls2csv (part of catdoc package)
        try:
            result = subprocess.run(
                ['xls2csv', xls_path],
                capture_output=True,
                text=True,
                check=True,
                timeout=120
            )
            text = result.stdout
            
            if text.strip():
                logger.info(f"Extracted {len(text)} characters from .xls file using xls2csv")
                return text
                
        except FileNotFoundError:
            logger.warning("xls2csv not found, trying textract")
        except subprocess.CalledProcessError as e:
            logger.warning(f"xls2csv failed with exit code {e.returncode}: {e.stderr}, trying textract")
        
        # Fallback to textract
        try:
            import textract
            text = textract.process(xls_path).decode('utf-8')
            
            if text.strip():
                logger.info(f"Extracted {len(text)} characters from .xls file using textract")
                return text
                
        except ImportError:
            logger.error("textract not installed")
        except Exception as e:
            logger.error(f"textract failed: {e}")
        
        raise ValueError(
            "Unable to extract text from .xls file. "
            "Please install catdoc (xls2csv) or textract library."
        )
        
    except Exception as e:
        logger.error(f"Failed to extract text from .xls file {xls_path}: {e}")
        raise


def is_legacy_office_file(file_path: str) -> Optional[str]:
    """
    Check if file is a legacy Office file that needs special processing.
    
    Args:
        file_path: Path to check
        
    Returns:
        File type ('ppt' or 'xls') if legacy format, None otherwise
    """
    lower_path = file_path.lower()
    
    if lower_path.endswith('.ppt') and not lower_path.endswith('.pptx'):
        return 'ppt'
    elif lower_path.endswith('.xls') and not lower_path.endswith('.xlsx'):
        return 'xls'
    
    return None